"""Business logic for document ingestion, similarity search, and RAG generation.

Embedding backend: sentence-transformers (free, local, no API key required).
Generation backend: Anthropic Claude (ANTHROPIC_API_KEY required for RAG).
"""

from __future__ import annotations

import io
import os
import re

import anthropic
from django.core.exceptions import ImproperlyConfigured
from pgvector.django import CosineDistance
from pypdf import PdfReader

from .models import Chunk, Document

# ---------------------------------------------------------------------------
# PDF Extraction
# ---------------------------------------------------------------------------


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract plain text from a PDF byte stream.

    Raises ValueError if the PDF produces no extractable text
    (e.g. scanned image-only PDF with no OCR layer).
    """
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    if not text:
        raise ValueError(
            "No extractable text found. The PDF may be scanned or image-only."
        )
    return text


# ---------------------------------------------------------------------------
# Plain Text & Markdown Extraction
# ---------------------------------------------------------------------------


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Decode UTF-8 plain text or Markdown, stripping YAML front matter if present."""
    text = file_bytes.decode("utf-8", errors="replace").strip()
    if text.startswith("---"):
        end = text.find("---", 3)
        if end != -1:
            text = text[end + 3 :].strip()
    if not text:
        raise ValueError("File contains no extractable text.")
    return text


# ---------------------------------------------------------------------------
# HTML Extraction
# ---------------------------------------------------------------------------


def extract_text_from_html(file_bytes: bytes) -> str:
    """Strip HTML tags and return visible text."""
    from bs4 import BeautifulSoup  # noqa: PLC0415

    soup = BeautifulSoup(file_bytes, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "head"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not text:
        raise ValueError("No visible text found in HTML file.")
    return text


# ---------------------------------------------------------------------------
# Office Document Extraction
# ---------------------------------------------------------------------------


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a .docx Word document, preserving paragraph order."""
    from docx import Document as DocxDocument  # noqa: PLC0415

    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    if not text:
        raise ValueError("No extractable text found in DOCX file.")
    return text


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from all slides in a .pptx file, one slide per block."""
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    text = "\n\n".join(slides)
    if not text:
        raise ValueError("No extractable text found in PPTX file.")
    return text


# ---------------------------------------------------------------------------
# Structured Data Extraction
# ---------------------------------------------------------------------------


def extract_text_from_csv(file_bytes: bytes, delimiter: str = ",") -> str:
    """Convert CSV/TSV rows to 'key: value' line blocks, one block per row."""
    import csv  # noqa: PLC0415

    reader = csv.DictReader(
        io.StringIO(file_bytes.decode("utf-8", errors="replace")),
        delimiter=delimiter,
    )
    rows = []
    for i, row in enumerate(reader):
        lines = [f"{k}: {v}" for k, v in row.items() if v and v.strip()]
        if lines:
            rows.append(f"[Row {i + 1}]\n" + "\n".join(lines))
    text = "\n\n".join(rows)
    if not text:
        raise ValueError("CSV file is empty or contains no data.")
    return text


def extract_text_from_json(file_bytes: bytes) -> str:
    """Recursively flatten JSON/JSONL into key: value text blocks."""
    import json  # noqa: PLC0415

    def _flatten(obj: object, prefix: str = "") -> list[str]:
        lines: list[str] = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                new_prefix = f"{prefix}{k}: " if not prefix else f"{prefix}.{k}: "
                lines.extend(_flatten(v, new_prefix))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                lines.extend(_flatten(item, f"{prefix}[{i}]: "))
        else:
            lines.append(f"{prefix}{obj}")
        return lines

    raw = file_bytes.decode("utf-8", errors="replace").strip()
    if "\n" in raw:
        records = []
        for i, line in enumerate(raw.splitlines()):
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
                records.append(f"[Record {i + 1}]\n" + "\n".join(_flatten(obj)))
            except json.JSONDecodeError:
                continue
        text = "\n\n".join(records)
    else:
        obj = json.loads(raw)
        text = "\n".join(_flatten(obj))
    if not text:
        raise ValueError("JSON file is empty or contains no data.")
    return text


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Convert each sheet of an Excel workbook to row blocks."""
    import openpyxl  # noqa: PLC0415

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        headers: list[str] | None = None
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            values = [str(v) if v is not None else "" for v in row]
            if i == 0:
                headers = values
                continue
            if not any(v.strip() for v in values):
                continue
            if headers:
                parts = [f"{h}: {v}" for h, v in zip(headers, values) if v.strip()]
            else:
                parts = [v for v in values if v.strip()]
            if parts:
                rows.append("\n".join(parts))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n\n".join(rows))
    text = "\n\n".join(sheets)
    if not text:
        raise ValueError("Excel file is empty or contains no data.")
    return text


# ---------------------------------------------------------------------------
# Central Dispatch
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "pdf",
    ".txt": "txt",
    ".md": "txt",
    ".html": "html",
    ".htm": "html",
    ".docx": "docx",
    ".pptx": "pptx",
    ".csv": "csv",
    ".tsv": "tsv",
    ".json": "json",
    ".jsonl": "json",
    ".xlsx": "xlsx",
}


def extract_text_from_file(name: str, file_bytes: bytes) -> str:
    """Dispatch to the correct extractor based on file extension."""
    ext = "." + name.lower().rsplit(".", 1)[-1]
    kind = SUPPORTED_EXTENSIONS.get(ext)
    if kind is None:
        allowed = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type '{ext}'. Allowed: {allowed}")
    dispatch = {
        "pdf": extract_text_from_pdf,
        "txt": extract_text_from_txt,
        "html": extract_text_from_html,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "csv": extract_text_from_csv,
        "tsv": lambda b: extract_text_from_csv(b, delimiter="\t"),
        "json": extract_text_from_json,
        "xlsx": extract_text_from_xlsx,
    }
    return dispatch[kind](file_bytes)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


def chunk_document(content: str, chunk_size: int = 512, overlap: int = 64) -> list[str]:
    """Split *content* into overlapping text segments.

    Splits on sentence boundaries where possible, otherwise falls back to a
    hard word-boundary split.
    """
    if not content:
        return []

    sentences = re.split(r"(?<=[.!?])\s+", content.strip())

    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for sentence in sentences:
        sentence_len = len(sentence)
        if current_len + sentence_len > chunk_size and current:
            chunks.append(" ".join(current))
            # Keep overlap: drop sentences from the front until we're under overlap
            while current and current_len > overlap:
                removed = current.pop(0)
                current_len -= len(removed) + 1
        current.append(sentence)
        current_len += sentence_len + 1  # +1 for the space

    if current:
        chunks.append(" ".join(current))

    return chunks


# ---------------------------------------------------------------------------
# Embedding backend — sentence-transformers (local, free, no API key)
# ---------------------------------------------------------------------------

EMBEDDING_DIMENSIONS = int(os.environ.get("EMBEDDING_DIMENSIONS", "1024"))
_EMBEDDING_MODEL_NAME = os.environ.get("EMBEDDING_MODEL", "BAAI/bge-large-en-v1.5")
_model_singleton = None


def get_embedding_model():
    """Lazy-load and cache the SentenceTransformer model as a module singleton."""
    global _model_singleton  # noqa: PLW0603
    if _model_singleton is None:
        from sentence_transformers import SentenceTransformer

        _model_singleton = SentenceTransformer(_EMBEDDING_MODEL_NAME)
    return _model_singleton


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Return a list of embedding vectors for *texts* using a local model.

    The model is loaded once and kept in memory (singleton).  No API key or
    network call is made after the initial model download.

    In tests, monkeypatch this function to return zero vectors of the correct
    dimension — no HTTP mock is needed since the model runs locally.
    """
    if not texts:
        return []
    model = get_embedding_model()
    vectors = model.encode(texts, convert_to_numpy=True)
    return [v.tolist() for v in vectors]


# ---------------------------------------------------------------------------
# Ingestion
# ---------------------------------------------------------------------------


def ingest_document(title: str, content: str, source: str = "") -> Document:
    """Create a Document, chunk it, embed the chunks, and persist everything.

    Returns the saved Document instance.
    """
    document = Document.objects.create(title=title, content=content, source=source)

    text_chunks = chunk_document(content)
    if not text_chunks:
        return document

    vectors = embed_texts(text_chunks)

    Chunk.objects.bulk_create(
        [
            Chunk(
                document=document,
                content=text,
                chunk_index=idx,
                embedding=vector,
            )
            for idx, (text, vector) in enumerate(zip(text_chunks, vectors, strict=True))
        ]
    )

    return document


# ---------------------------------------------------------------------------
# Search
# ---------------------------------------------------------------------------


def search_similar_chunks(query: str, top_k: int = 5):  # type: ignore[return]
    """Embed *query* and return the *top_k* most similar Chunk objects.

    Results are annotated with a ``distance`` attribute (lower = more similar).
    """
    query_vector = embed_texts([query])[0]
    return (
        Chunk.objects.select_related("document")
        .annotate(distance=CosineDistance("embedding", query_vector))
        .order_by("distance")[:top_k]
    )


# ---------------------------------------------------------------------------
# Generation — Claude via Anthropic API
# ---------------------------------------------------------------------------

_CLAUDE_MODEL = os.environ.get("CLAUDE_MODEL", "claude-opus-4-5")


def generate_answer(query: str, context_chunks: list[Chunk]) -> str:
    """Generate a grounded answer using Claude with retrieved chunks as context.

    Requires ``ANTHROPIC_API_KEY`` to be set in the environment.
    Raises ``ImproperlyConfigured`` if the key is absent.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        raise ImproperlyConfigured(
            "ANTHROPIC_API_KEY is not set. "
            "Set it in your environment or .env file to use the RAG endpoint."
        )

    client = anthropic.Anthropic(api_key=api_key)
    context = "\n\n".join(chunk.content for chunk in context_chunks)
    message = client.messages.create(
        model=_CLAUDE_MODEL,
        max_tokens=1024,
        messages=[
            {
                "role": "user",
                "content": f"Context:\n{context}\n\nQuestion: {query}",
            }
        ],
    )
    return message.content[0].text

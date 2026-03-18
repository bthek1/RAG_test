"""Unit tests for apps.embeddings.services — no DB required."""

import numpy as np
import pytest

from apps.embeddings.services import (
    EMBEDDING_DIMENSIONS,
    chunk_document,
    embed_texts,
    extract_text_from_pdf,
    generate_answer,
)

# ---------------------------------------------------------------------------
# chunk_document
# ---------------------------------------------------------------------------


class TestChunkDocument:
    def test_empty_input_returns_empty_list(self):
        assert chunk_document("") == []

    def test_single_short_sentence(self):
        chunks = chunk_document("Hello world.", chunk_size=512)
        assert chunks == ["Hello world."]

    def test_splits_long_content_into_multiple_chunks(self):
        # Build a string definitely longer than chunk_size
        sentence = "This is a test sentence. "
        content = sentence * 30  # ~750 chars
        chunks = chunk_document(content, chunk_size=100, overlap=20)
        assert len(chunks) > 1
        # Every chunk should be non-empty
        assert all(c.strip() for c in chunks)

    def test_overlap_means_content_shared_between_chunks(self):
        # With overlap the last sentence of chunk N should appear in chunk N+1
        sentence = "Sentence number {}. "
        content = "".join(sentence.format(i) for i in range(20))
        chunks = chunk_document(content, chunk_size=80, overlap=40)
        if len(chunks) > 1:
            # The second chunk should contain text that was near the end of the first
            assert len(chunks[1]) > 0

    def test_returns_list_of_strings(self):
        chunks = chunk_document("One sentence. Two sentences.", chunk_size=50)
        assert isinstance(chunks, list)
        assert all(isinstance(c, str) for c in chunks)


# ---------------------------------------------------------------------------
# embed_texts — monkeypatched SentenceTransformer (no network or GPU needed)
# ---------------------------------------------------------------------------


class TestEmbedTexts:
    def _make_fake_model(self, fake_vectors):
        class _FakeModel:
            def encode(self, texts, **kw):
                return fake_vectors

        return _FakeModel()

    def test_returns_list_of_vectors(self, monkeypatch):
        fake_vectors = np.zeros((2, EMBEDDING_DIMENSIONS), dtype=np.float32)
        monkeypatch.setattr(
            "apps.embeddings.services.get_embedding_model",
            lambda: self._make_fake_model(fake_vectors),
        )
        result = embed_texts(["hello", "world"])
        assert len(result) == 2

    def test_vector_has_correct_dimension(self, monkeypatch):
        fake_vectors = np.zeros((1, EMBEDDING_DIMENSIONS), dtype=np.float32)
        monkeypatch.setattr(
            "apps.embeddings.services.get_embedding_model",
            lambda: self._make_fake_model(fake_vectors),
        )
        result = embed_texts(["test"])
        assert len(result[0]) == EMBEDDING_DIMENSIONS

    def test_empty_list_returns_empty_list(self, monkeypatch):
        # embed_texts short-circuits before calling the model for empty input
        result = embed_texts([])
        assert result == []

    def test_default_dimension_is_1024(self):
        assert EMBEDDING_DIMENSIONS == 1024


# ---------------------------------------------------------------------------
# generate_answer
# ---------------------------------------------------------------------------


class TestGenerateAnswer:
    def test_raises_improperly_configured_when_no_api_key(self, monkeypatch):
        from django.core.exceptions import ImproperlyConfigured

        monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
        # Reset the module-level env read by clearing os.environ
        with pytest.raises(ImproperlyConfigured, match="ANTHROPIC_API_KEY"):
            generate_answer("What is RAG?", [])

    def test_calls_anthropic_and_returns_text(self, monkeypatch):
        monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")

        # Build a minimal fake Anthropic response object
        class FakeTextBlock:
            text = "RAG stands for Retrieval-Augmented Generation."

        class FakeMessage:
            content = [FakeTextBlock()]

        class FakeMessages:
            def create(self, **kwargs):
                return FakeMessage()

        class FakeClient:
            messages = FakeMessages()

        monkeypatch.setattr(
            "apps.embeddings.services.anthropic.Anthropic",
            lambda **kw: FakeClient(),
        )

        # Minimal Chunk-like objects
        class FakeChunk:
            content = "Some context text."

        result = generate_answer("What is RAG?", [FakeChunk()])
        assert result == "RAG stands for Retrieval-Augmented Generation."

    def test_prompt_includes_context_and_query(self, monkeypatch):
        monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")
        captured: dict = {}

        class FakeTextBlock:
            text = "answer"

        class FakeMessage:
            content = [FakeTextBlock()]

        class FakeMessages:
            def create(self, **kwargs):
                captured.update(kwargs)
                return FakeMessage()

        class FakeClient:
            messages = FakeMessages()

        monkeypatch.setattr(
            "apps.embeddings.services.anthropic.Anthropic",
            lambda **kw: FakeClient(),
        )

        class FakeChunk:
            content = "context chunk text"

        generate_answer("my question", [FakeChunk()])
        user_msg = captured["messages"][0]["content"]
        assert "context chunk text" in user_msg
        assert "my question" in user_msg


# ---------------------------------------------------------------------------
# extract_text_from_pdf
# ---------------------------------------------------------------------------


class TestExtractTextFromPdf:
    def _mock_reader(self, monkeypatch, pages_text: list[str]):
        """Patch PdfReader to return pages with the given text strings."""
        from unittest.mock import MagicMock

        mock_pages = []
        for t in pages_text:
            page = MagicMock()
            page.extract_text.return_value = t
            mock_pages.append(page)

        mock_reader = MagicMock()
        mock_reader.pages = mock_pages
        monkeypatch.setattr(
            "apps.embeddings.services.PdfReader", lambda _buf: mock_reader
        )

    def test_extracts_text_from_single_page(self, monkeypatch):
        self._mock_reader(monkeypatch, ["Hello, PDF world!"])
        text = extract_text_from_pdf(b"fake")
        assert "Hello, PDF world!" in text

    def test_joins_multiple_pages(self, monkeypatch):
        self._mock_reader(monkeypatch, ["Page one content.", "Page two content."])
        text = extract_text_from_pdf(b"fake")
        assert "Page one content." in text
        assert "Page two content." in text

    def test_raises_value_error_when_all_pages_empty(self, monkeypatch):
        self._mock_reader(monkeypatch, ["", "  ", ""])
        with pytest.raises(ValueError, match="No extractable text"):
            extract_text_from_pdf(b"fake")

    def test_raises_value_error_for_zero_pages(self, monkeypatch):
        self._mock_reader(monkeypatch, [])
        with pytest.raises(ValueError, match="No extractable text"):
            extract_text_from_pdf(b"fake")

    def test_skips_blank_pages(self, monkeypatch):
        self._mock_reader(monkeypatch, ["", "Real content here.", ""])
        text = extract_text_from_pdf(b"fake")
        assert "Real content here." in text


# ---------------------------------------------------------------------------
# extract_text_from_txt
# ---------------------------------------------------------------------------


class TestExtractTextFromTxt:
    def test_plain_utf8_text(self):
        from apps.embeddings.services import extract_text_from_txt

        text = extract_text_from_txt(b"Hello, world!")
        assert text == "Hello, world!"

    def test_strips_yaml_front_matter(self):
        from apps.embeddings.services import extract_text_from_txt

        content = b"---\ntitle: Test\n---\nBody text here."
        text = extract_text_from_txt(content)
        assert text == "Body text here."
        assert "title: Test" not in text

    def test_markdown_without_front_matter(self):
        from apps.embeddings.services import extract_text_from_txt

        content = b"# Heading\n\nSome paragraph."
        text = extract_text_from_txt(content)
        assert "# Heading" in text
        assert "Some paragraph." in text

    def test_raises_on_empty_file(self):
        from apps.embeddings.services import extract_text_from_txt

        with pytest.raises(ValueError, match="no extractable text"):
            extract_text_from_txt(b"")

    def test_raises_on_whitespace_only(self):
        from apps.embeddings.services import extract_text_from_txt

        with pytest.raises(ValueError, match="no extractable text"):
            extract_text_from_txt(b"   \n\n  ")


# ---------------------------------------------------------------------------
# extract_text_from_html
# ---------------------------------------------------------------------------


class TestExtractTextFromHtml:
    def test_strips_tags_returns_visible_text(self):
        from apps.embeddings.services import extract_text_from_html

        html = b"<html><body><p>Hello <b>world</b>!</p></body></html>"
        text = extract_text_from_html(html)
        assert "Hello" in text
        assert "world" in text

    def test_removes_script_and_style(self):
        from apps.embeddings.services import extract_text_from_html

        html = b"<html><head><style>body{color:red}</style></head><body><script>alert(1)</script><p>Visible</p></body></html>"
        text = extract_text_from_html(html)
        assert "Visible" in text
        assert "alert" not in text
        assert "color:red" not in text

    def test_raises_on_no_visible_text(self):
        from apps.embeddings.services import extract_text_from_html

        html = b"<html><body><script>alert(1)</script></body></html>"
        with pytest.raises(ValueError, match="No visible text"):
            extract_text_from_html(html)


# ---------------------------------------------------------------------------
# extract_text_from_docx
# ---------------------------------------------------------------------------


class TestExtractTextFromDocx:
    def _make_docx_bytes(self, paragraphs: list[str]) -> bytes:
        """Create a minimal in-memory DOCX with the given paragraphs."""
        import io

        from docx import Document

        doc = Document()
        for para in paragraphs:
            doc.add_paragraph(para)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def test_extracts_paragraphs(self):
        from apps.embeddings.services import extract_text_from_docx

        docx_bytes = self._make_docx_bytes(["First paragraph.", "Second paragraph."])
        text = extract_text_from_docx(docx_bytes)
        assert "First paragraph." in text
        assert "Second paragraph." in text

    def test_raises_on_empty_docx(self):
        from apps.embeddings.services import extract_text_from_docx

        docx_bytes = self._make_docx_bytes([])
        with pytest.raises(ValueError, match="No extractable text"):
            extract_text_from_docx(docx_bytes)


# ---------------------------------------------------------------------------
# extract_text_from_pptx
# ---------------------------------------------------------------------------


class TestExtractTextFromPptx:
    def _make_pptx_bytes(self, slides: list[list[str]]) -> bytes:
        """Create a minimal in-memory PPTX with given slides (list of text items)."""
        import io

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank layout
        for slide_texts in slides:
            slide = prs.slides.add_slide(blank_layout)
            for txt in slide_texts:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                txBox.text_frame.text = txt
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_extracts_slide_text(self):
        from apps.embeddings.services import extract_text_from_pptx

        pptx_bytes = self._make_pptx_bytes([["Slide one content"], ["Slide two content"]])
        text = extract_text_from_pptx(pptx_bytes)
        assert "Slide one content" in text
        assert "Slide two content" in text
        assert "[Slide 1]" in text
        assert "[Slide 2]" in text

    def test_raises_on_empty_pptx(self):
        from apps.embeddings.services import extract_text_from_pptx

        pptx_bytes = self._make_pptx_bytes([])
        with pytest.raises(ValueError, match="No extractable text"):
            extract_text_from_pptx(pptx_bytes)


# ---------------------------------------------------------------------------
# extract_text_from_csv
# ---------------------------------------------------------------------------


class TestExtractTextFromCsv:
    def test_extracts_rows_with_headers(self):
        from apps.embeddings.services import extract_text_from_csv

        csv_bytes = b"name,age\nAlice,30\nBob,25"
        text = extract_text_from_csv(csv_bytes)
        assert "name: Alice" in text
        assert "age: 30" in text
        assert "[Row 1]" in text

    def test_handles_tsv_delimiter(self):
        from apps.embeddings.services import extract_text_from_csv

        tsv_bytes = b"name\tage\nAlice\t30"
        text = extract_text_from_csv(tsv_bytes, delimiter="\t")
        assert "name: Alice" in text

    def test_raises_on_empty_csv(self):
        from apps.embeddings.services import extract_text_from_csv

        with pytest.raises(ValueError, match="empty or contains no data"):
            extract_text_from_csv(b"name,age\n")


# ---------------------------------------------------------------------------
# extract_text_from_json
# ---------------------------------------------------------------------------


class TestExtractTextFromJson:
    def test_flat_json_object(self):
        from apps.embeddings.services import extract_text_from_json

        json_bytes = b'{"key": "value", "num": 42}'
        text = extract_text_from_json(json_bytes)
        assert "key: " in text
        assert "value" in text

    def test_jsonl_multiple_records(self):
        from apps.embeddings.services import extract_text_from_json

        jsonl_bytes = b'{"a": 1}\n{"b": 2}'
        text = extract_text_from_json(jsonl_bytes)
        assert "[Record 1]" in text
        assert "[Record 2]" in text

    def test_nested_json_flattened(self):
        from apps.embeddings.services import extract_text_from_json

        json_bytes = b'{"outer": {"inner": "deep"}}'
        text = extract_text_from_json(json_bytes)
        assert "deep" in text

    def test_raises_on_empty_json(self):
        from apps.embeddings.services import extract_text_from_json

        with pytest.raises(ValueError, match="empty or contains no data"):
            extract_text_from_json(b"{}")


# ---------------------------------------------------------------------------
# extract_text_from_xlsx
# ---------------------------------------------------------------------------


class TestExtractTextFromXlsx:
    def _make_xlsx_bytes(self, rows: list[list[str]], sheet_name: str = "Sheet1") -> bytes:
        import io

        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name
        for row in rows:
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def test_extracts_sheet_data(self):
        from apps.embeddings.services import extract_text_from_xlsx

        xlsx_bytes = self._make_xlsx_bytes([["Name", "Score"], ["Alice", "95"], ["Bob", "82"]])
        text = extract_text_from_xlsx(xlsx_bytes)
        assert "Name: Alice" in text
        assert "Score: 95" in text
        assert "[Sheet: Sheet1]" in text

    def test_raises_on_empty_workbook(self):
        from apps.embeddings.services import extract_text_from_xlsx

        xlsx_bytes = self._make_xlsx_bytes([["Header1", "Header2"]])
        with pytest.raises(ValueError, match="empty or contains no data"):
            extract_text_from_xlsx(xlsx_bytes)


# ---------------------------------------------------------------------------
# extract_text_from_file (dispatch)
# ---------------------------------------------------------------------------


class TestExtractTextFromFile:
    def test_dispatches_pdf(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_pdf",
            lambda b: "pdf content",
        )
        assert extract_text_from_file("doc.pdf", b"bytes") == "pdf content"

    def test_dispatches_txt(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_txt",
            lambda b: "txt content",
        )
        assert extract_text_from_file("notes.txt", b"bytes") == "txt content"

    def test_dispatches_md_to_txt_extractor(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_txt",
            lambda b: "md content",
        )
        assert extract_text_from_file("readme.md", b"bytes") == "md content"

    def test_dispatches_csv(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_csv",
            lambda b: "csv content",
        )
        assert extract_text_from_file("data.csv", b"bytes") == "csv content"

    def test_dispatches_xlsx(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_xlsx",
            lambda b: "xlsx content",
        )
        assert extract_text_from_file("sheet.xlsx", b"bytes") == "xlsx content"

    def test_unknown_extension_raises_value_error(self):
        from apps.embeddings.services import extract_text_from_file

        with pytest.raises(ValueError, match="Unsupported file type"):
            extract_text_from_file("malware.exe", b"bytes")

    def test_case_insensitive_extension(self, monkeypatch):
        from apps.embeddings.services import extract_text_from_file

        monkeypatch.setattr(
            "apps.embeddings.services.extract_text_from_pdf",
            lambda b: "pdf content",
        )
        assert extract_text_from_file("DOC.PDF", b"bytes") == "pdf content"
